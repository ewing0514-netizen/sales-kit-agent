"""
generator.py - Sales Kit 核心生成邏輯
1. 呼叫 Claude API 生成各頁文案 (JSON)
2. 使用 python-pptx 組裝成 .pptx 檔案
"""

import os
import json
import re
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ══════════════════════════════════════════
# 全域色彩定義（商務深藍風格）
# ══════════════════════════════════════════
C_PRIMARY   = RGBColor(0x1B, 0x3A, 0x6B)   # 深藍（主色）
C_ACCENT    = RGBColor(0x2D, 0x6B, 0xE8)   # 亮藍（強調）
C_ACCENT2   = RGBColor(0x06, 0x64, 0xC8)   # 中藍
C_SOFT_BLUE = RGBColor(0x25, 0x50, 0x9E)   # 半透明裝飾藍
C_LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFB)   # 淺灰背景
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)    # 白色
C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)   # 深色文字
C_GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)   # 灰色次要文字
C_GOLD      = RGBColor(0xF5, 0xA6, 0x23)   # 金色點綴
C_PALE_BLUE = RGBColor(0xAD, 0xC8, 0xFF)   # 淡藍文字

# 投影片尺寸（16:9 寬螢幕，英吋）
SW = 13.33  # slide width
SH = 7.5    # slide height


# ══════════════════════════════════════════
# Claude API 相關函式
# ══════════════════════════════════════════

def get_kit_display_name(kit_type: str) -> str:
    """將 kit_type 代碼轉換為中文顯示名稱"""
    mapping = {
        'bni':     'BNI 主題簡報｜LCD 333 架構',
        'product': '產品與服務銷售提案',
        'brand':   '品牌合作提案',
        'event':   '活動推廣 Sales Kit',
    }
    return mapping.get(kit_type, kit_type)


def build_prompt(form_data: dict) -> str:
    """根據表單資料組裝 Claude API 提示詞"""
    kit_type      = form_data.get('kit_type', 'product')
    company_name  = form_data.get('company_name', '')
    core_value    = form_data.get('core_value', '')
    target        = form_data.get('target_audience', '')
    data_hl       = form_data.get('data_highlights', '')
    contact       = form_data.get('contact_info', '')
    extra         = form_data.get('extra_info', '') or '（無補充）'
    display_name  = get_kit_display_name(kit_type)

    # 各類型對應的投影片結構
    structures = {
        'bni': [
            # ── 介紹 3頁（自我介紹）──
            'cover    → 封面｜開場（姓名、公司、產業別）',
            'intro    → 我在做什麼？（產品/服務/年資/專業背景）',
            'intro    → 為什麼做這事業？（初衷故事、使命感、改變點）',
            # ── 概要 3頁（成功案例）──
            'case     → 成功案例 1（客戶問題 → 解決方案 → 具體成果）',
            'case     → 成功案例 2（客戶問題 → 解決方案 → 具體成果）',
            'case     → 成功案例 3（客戶問題 → 解決方案 → 具體成果）',
            # ── 內容 3頁（如何滿足需求）──
            'steps    → 希望引薦的客戶（一般引薦、理想引薦、夢幻引薦）',
            'intro    → 總結｜最重要的一件事（核心訴求精華）',
            'cta      → Slogan｜記憶點口號與聯絡資訊',
        ],
        'product': [
            'cover    → 封面（公司 + 產品名稱）',
            'problem  → 您面臨的挑戰',
            'solution → 我們的解決方案',
            'features → 三大核心功能',
            'stats    → 成果數據',
            'intro    → 成功案例',
            'steps    → 合作流程（三步驟）',
            'cta      → 立即開始',
        ],
        'brand': [
            'cover    → 封面（品牌名稱）',
            'intro    → 品牌故事與願景',
            'problem  → 市場機會',
            'solution → 我們的定位',
            'features → 三大品牌優勢',
            'stats    → 品牌數據',
            'steps    → 合作模式',
            'cta      → 合作提案',
        ],
        'event': [
            'cover    → 封面（活動名稱）',
            'intro    → 活動概覽',
            'problem  → 為什麼現在需要這個活動',
            'features → 活動三大亮點',
            'stats    → 活動規模數據',
            'solution → 精彩內容與議程',
            'steps    → 如何報名（三步驟）',
            'cta      → 立即報名',
        ],
    }

    slide_count = len(structures.get(kit_type, structures['product']))
    slide_list = '\n'.join(
        [f'  {i+1}. {s}' for i, s in enumerate(structures.get(kit_type, structures['product']))]
    )

    # BNI 專屬附加說明
    bni_extra = ''
    if kit_type == 'bni':
        bni_extra = """
【BNI LCD 333 架構說明】
這是 BNI 標準主題簡報架構，分三大區塊：
- 介紹（3頁）：讓成員認識你是誰、做什麼、為何而做
- 概要（3頁）：用3個真實成功案例展示你的能力
- 內容（3頁）：說明理想引薦對象、總結核心訊息、留下記憶點

成功案例重點：要具體、有數字、說出客戶的轉變（before → after）
希望引薦投影片：明確說出一般/理想/夢幻三種引薦對象
"""

    return f"""你是台灣頂尖的商業簡報文案專家，擅長撰寫具說服力的銷售提案。

請根據以下客戶資料，為「{display_name}」生成一份 PowerPoint 簡報的完整文案。
{bni_extra}
【客戶資訊】
- 姓名／公司／品牌：{company_name}
- 核心訴求：{core_value}
- 目標受眾：{target}
- 數據亮點：{data_hl}
- 聯絡資訊：{contact}
- 補充說明：{extra}

【需要生成的投影片（共 {slide_count} 頁）】
{slide_list}

【JSON 輸出格式說明】
根據投影片類型填入對應欄位：

cover 類型：
  {{"type":"cover","title":"主標題","subtitle":"副標題說明","tagline":"記憶點口號"}}

intro / problem / solution 類型：
  {{"type":"intro","title":"標題","content":"50字以內的引言說明","points":["要點1（20字內）","要點2","要點3"]}}

case 類型（成功案例，BNI 專用）：
  {{"type":"case","title":"成功案例 X｜客戶類型描述","client":"客戶背景一行描述","problem":"客戶原本面臨的問題或痛點（40字內）","solution":"你提供的解決方案（40字內）","result":"具體成果或數字化結果（40字內）"}}

features 類型（固定 3 個 items）：
  {{"type":"features","title":"標題","items":[{{"title":"特色名稱","description":"30字說明"}},{{"title":"...","description":"..."}},{{"title":"...","description":"..."}}]}}

stats 類型（固定 3 個 items）：
  {{"type":"stats","title":"標題","items":[{{"number":"數字/比率","label":"指標名稱","description":"補充說明"}},{{"number":"...","label":"...","description":"..."}},{{"number":"...","label":"...","description":"..."}}]}}

steps 類型（BNI 用於希望引薦）：
  {{"type":"steps","title":"標題","items":[{{"step":"一般","title":"一般引薦對象","description":"具體描述此類客戶特徵"}},{{"step":"理想","title":"理想引薦對象","description":"具體描述此類客戶特徵"}},{{"step":"夢幻","title":"夢幻引薦對象","description":"具體描述理想中最完美的客戶"}}]}}

cta 類型：
  {{"type":"cta","title":"簡短有力的 Slogan（10字內）","subtitle":"一句話說明你能帶來的價值","contact":"{contact}"}}

【重要規則】
1. 全程繁體中文
2. 文案真實、具體、有說服力，符合台灣商業文化
3. 成功案例必須有 before/after 的轉變感，有具體數字更好
4. 每個要點控制在 20 字以內
5. 直接輸出 JSON，不要有任何說明文字

輸出：
{{"slides": [ ...{slide_count} 個物件... ]}}"""


def call_claude(form_data: dict) -> dict:
    """呼叫 Claude API，取得投影片 JSON 資料"""
    api_key = os.environ.get('ANTHROPIC_API_KEY')
    if not api_key:
        raise ValueError("未設定 ANTHROPIC_API_KEY 環境變數")

    client = anthropic.Anthropic(api_key=api_key)
    prompt = build_prompt(form_data)

    message = client.messages.create(
        model='claude-sonnet-4-6',
        max_tokens=4096,
        messages=[{'role': 'user', 'content': prompt}],
    )

    raw = message.content[0].text
    return parse_json(raw)


def parse_json(text: str) -> dict:
    """從 Claude 回應文字中解析出 JSON"""
    text = text.strip()

    # 嘗試直接解析
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # 嘗試擷取 markdown 程式碼區塊
    m = re.search(r'```(?:json)?\s*([\s\S]*?)```', text)
    if m:
        try:
            return json.loads(m.group(1).strip())
        except json.JSONDecodeError:
            pass

    # 嘗試擷取最外層 {...}
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end > start:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass

    raise ValueError("無法解析 Claude 回傳的 JSON，請再試一次")


# ══════════════════════════════════════════
# PPTX 輔助函式
# ══════════════════════════════════════════

def blank_slide(prs: Presentation):
    """新增空白投影片（使用 Blank 版面）"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor):
    """設定投影片背景色"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color: RGBColor):
    """新增無邊框矩形色塊（單位：英吋）"""
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def txt(slide, text: str, x, y, w, h,
        size=18, bold=False, italic=False,
        color: RGBColor = C_WHITE,
        align=PP_ALIGN.LEFT):
    """新增單段文字框（單位：英吋）"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullets(slide, points: list, x, y, w, h,
            size=16, color: RGBColor = C_DARK_TEXT, prefix='▶  '):
    """新增項目清單文字框"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f'{prefix}{point}'
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_logo(slide, logo_path: str, x, y, w=2.2, h=1.4):
    """嘗試插入 Logo 圖片，失敗時靜默跳過"""
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y),
                                     width=Inches(w), height=Inches(h))
        except Exception:
            pass  # 圖片格式不支援時忽略


# ══════════════════════════════════════════
# 各類型投影片建立函式
# ══════════════════════════════════════════

def slide_cover(prs, data: dict, logo_path):
    """封面投影片：深藍背景 + 金色裝飾線 + 公司名稱"""
    s = blank_slide(prs)
    set_bg(s, C_PRIMARY)

    # 頂部亮藍裝飾條
    rect(s, 0, 0, SW, 0.08, C_ACCENT)
    # 底部金色裝飾條
    rect(s, 0, SH - 0.08, SW, 0.08, C_GOLD)
    # 右側半透明裝飾塊
    rect(s, 9.3, 0, SW - 9.3, SH, C_SOFT_BLUE)

    # 左側金色裝飾豎線
    rect(s, 0.7, 1.8, 0.07, 4.0, C_GOLD)

    title    = data.get('title', '')
    subtitle = data.get('subtitle', '')
    tagline  = data.get('tagline', '')

    # 主標題（公司名稱）
    txt(s, title, 1.0, 1.9, 7.8, 1.5, size=38, bold=True, color=C_WHITE)
    # 副標題
    if subtitle:
        txt(s, subtitle, 1.0, 3.5, 7.8, 0.9, size=21, color=C_PALE_BLUE)
    # 口號
    if tagline:
        txt(s, f'「{tagline}」', 1.0, 4.5, 7.8, 0.7, size=15, italic=True,
            color=RGBColor(0x9B, 0xB3, 0xE8))

    # Logo
    add_logo(s, logo_path, x=9.7, y=0.4)

    # 右下角標籤
    txt(s, 'SALES KIT', 10.0, 6.8, 3.0, 0.5, size=10, bold=True,
        color=RGBColor(0x4A, 0x7A, 0xD4), align=PP_ALIGN.RIGHT)


def slide_content(prs, data: dict):
    """通用內容投影片：淺灰背景 + 深藍標題列 + 項目清單"""
    s = blank_slide(prs)
    set_bg(s, C_LIGHT_BG)

    # 頂部深藍標題列
    rect(s, 0, 0, SW, 1.3, C_PRIMARY)
    # 標題列右上角亮藍小塊
    rect(s, SW - 1.5, 0, 1.5, 1.3, C_ACCENT)

    title   = data.get('title', '')
    content = data.get('content', '')
    points  = data.get('points', [])

    txt(s, title, 0.6, 0.22, 11.5, 0.9, size=26, bold=True, color=C_WHITE)

    # 左側亮藍裝飾豎線
    rect(s, 0.35, 1.45, 0.07, 5.7, C_ACCENT)

    # 引言說明
    if content:
        txt(s, content, 0.65, 1.5, 12.2, 0.85, size=16, color=C_GRAY_TEXT)

    # 要點清單
    if points:
        y_start = 2.5 if content else 1.7
        bullets(s, points, 0.7, y_start, 12.0, 4.5, size=18)

    # 底部亮藍裝飾條
    rect(s, 0, SH - 0.1, SW, 0.1, C_ACCENT)


def slide_features(prs, data: dict):
    """特色展示投影片：3 張深色卡片"""
    s = blank_slide(prs)
    set_bg(s, C_LIGHT_BG)

    rect(s, 0, 0, SW, 1.3, C_PRIMARY)
    rect(s, SW - 1.5, 0, 1.5, 1.3, C_ACCENT)
    txt(s, data.get('title', '核心特色'), 0.6, 0.22, 11.5, 0.9,
        size=26, bold=True, color=C_WHITE)

    items  = (data.get('items') or [])[:3]
    cw     = 3.9
    cy     = 1.55
    ch     = 5.1
    colors = [C_PRIMARY, C_ACCENT2, C_ACCENT]

    for i, item in enumerate(items):
        cx = 0.37 + i * (cw + 0.37)
        rect(s, cx, cy, cw, ch, colors[i % 3])

        # 卡片頂部金色線
        rect(s, cx, cy, cw, 0.08, C_GOLD)

        # 序號
        txt(s, f'0{i + 1}', cx + 0.25, cy + 0.22, 1.2, 0.7,
            size=30, bold=True, color=C_GOLD)

        # 特色標題
        txt(s, item.get('title', ''), cx + 0.25, cy + 1.05, cw - 0.5, 0.65,
            size=17, bold=True, color=C_WHITE)

        # 特色說明
        txt(s, item.get('description', ''), cx + 0.25, cy + 1.85, cw - 0.5, 2.8,
            size=13, color=C_PALE_BLUE)

    rect(s, 0, SH - 0.1, SW, 0.1, C_ACCENT)


def slide_stats(prs, data: dict):
    """數據展示投影片：3 個大數字區塊"""
    s = blank_slide(prs)
    set_bg(s, C_WHITE)

    rect(s, 0, 0, SW, 1.3, C_PRIMARY)
    rect(s, SW - 1.5, 0, 1.5, 1.3, C_ACCENT)
    txt(s, data.get('title', '數據亮點'), 0.6, 0.22, 11.5, 0.9,
        size=26, bold=True, color=C_WHITE)

    items = (data.get('items') or [])[:3]
    bw    = 3.9
    by    = 1.65

    for i, item in enumerate(items):
        bx = 0.37 + i * (bw + 0.37)

        # 數據框
        rect(s, bx, by, bw, 5.1, C_LIGHT_BG)
        # 頂部強調線
        rect(s, bx, by, bw, 0.1, C_ACCENT)

        # 大數字
        txt(s, item.get('number', ''), bx + 0.15, by + 0.35, bw - 0.3, 1.5,
            size=46, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

        # 指標名稱
        txt(s, item.get('label', ''), bx + 0.15, by + 1.95, bw - 0.3, 0.65,
            size=15, bold=True, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

        # 補充說明
        txt(s, item.get('description', ''), bx + 0.15, by + 2.75, bw - 0.3, 1.9,
            size=13, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

    rect(s, 0, SH - 0.1, SW, 0.1, C_ACCENT)


def slide_steps(prs, data: dict):
    """步驟說明投影片：3 個帶編號的步驟"""
    s = blank_slide(prs)
    set_bg(s, C_LIGHT_BG)

    rect(s, 0, 0, SW, 1.3, C_PRIMARY)
    rect(s, SW - 1.5, 0, 1.5, 1.3, C_ACCENT)
    txt(s, data.get('title', '如何開始'), 0.6, 0.22, 11.5, 0.9,
        size=26, bold=True, color=C_WHITE)

    items = (data.get('items') or [])[:3]
    sw_   = 3.8   # step width
    sy    = 1.65

    # 步驟間連接橫線
    if len(items) > 1:
        rect(s, 2.8, sy + 0.45, 7.5, 0.07, C_ACCENT)

    for i, item in enumerate(items):
        sx = 0.5 + i * (sw_ + 0.6)

        # 步驟圓形背景（用正方形代替）
        circle_x = sx + sw_ / 2 - 0.52
        rect(s, circle_x, sy, 1.04, 1.04, C_PRIMARY)
        # 圓形金色邊框模擬
        rect(s, circle_x - 0.06, sy - 0.06, 1.16, 1.16, C_GOLD)
        rect(s, circle_x, sy, 1.04, 1.04, C_PRIMARY)

        txt(s, item.get('step', f'0{i+1}'), circle_x, sy + 0.12, 1.04, 0.8,
            size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 步驟標題
        txt(s, item.get('title', ''), sx, sy + 1.2, sw_, 0.65,
            size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

        # 步驟說明
        txt(s, item.get('description', ''), sx, sy + 1.95, sw_, 3.4,
            size=13, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    rect(s, 0, SH - 0.1, SW, 0.1, C_ACCENT)


def slide_case(prs, data: dict):
    """成功案例投影片（BNI 333架構）：3欄式 問題｜解決方案｜成果"""
    s = blank_slide(prs)
    set_bg(s, C_LIGHT_BG)

    # 頂部深藍標題列
    rect(s, 0, 0, SW, 1.3, C_PRIMARY)
    rect(s, SW - 1.5, 0, 1.5, 1.3, C_ACCENT)
    txt(s, data.get('title', '成功案例'), 0.6, 0.22, 11.5, 0.9,
        size=24, bold=True, color=C_WHITE)

    # 客戶背景標籤
    client = data.get('client', '')
    if client:
        txt(s, f'客戶背景：{client}', 0.5, 1.32, SW - 1.0, 0.42,
            size=12, italic=True, color=C_GRAY_TEXT)

    # 三欄顏色定義
    col_data = [
        ('問題／挑戰',  data.get('problem', ''),  RGBColor(0xC0, 0x39, 0x2B), RGBColor(0xFF, 0xF0, 0xEE)),
        ('解決方案',    data.get('solution', ''), C_ACCENT,                   RGBColor(0xEE, 0xF4, 0xFF)),
        ('成果數據',    data.get('result', ''),   RGBColor(0x16, 0xA3, 0x4A), RGBColor(0xEE, 0xFB, 0xF3)),
    ]

    cw = 3.9
    cy = 1.82
    ch = 5.3

    for i, (label, content, header_color, bg_color) in enumerate(col_data):
        cx = 0.37 + i * (cw + 0.37)

        # 欄位底框（淡色）
        rect(s, cx, cy, cw, ch, bg_color)
        # 欄位頂部彩色標題列
        rect(s, cx, cy, cw, 0.65, header_color)

        # 欄位標題
        txt(s, label, cx + 0.15, cy + 0.1, cw - 0.3, 0.5,
            size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 欄位內容
        txt(s, content, cx + 0.2, cy + 0.82, cw - 0.4, ch - 1.1,
            size=15, color=C_DARK_TEXT)

    rect(s, 0, SH - 0.1, SW, 0.1, C_ACCENT)


def slide_cta(prs, data: dict, logo_path):
    """行動呼籲投影片：深藍背景 + 白色文字 + 聯絡資訊"""
    s = blank_slide(prs)
    set_bg(s, C_PRIMARY)

    rect(s, 0, 0, SW, 0.08, C_ACCENT)
    rect(s, 0, SH - 0.08, SW, 0.08, C_GOLD)

    title    = data.get('title', '立即行動')
    subtitle = data.get('subtitle', '')
    contact  = data.get('contact', '')

    # 中央裝飾線
    rect(s, 3.5, 2.25, 6.3, 0.07, C_ACCENT)

    txt(s, title, 1.2, 1.3, 10.9, 1.2,
        size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    if subtitle:
        txt(s, subtitle, 1.2, 2.6, 10.9, 0.95,
            size=19, color=C_PALE_BLUE, align=PP_ALIGN.CENTER)

    if contact:
        # 聯絡資訊底框
        rect(s, 3.0, 3.8, 7.3, 1.4, RGBColor(0x25, 0x50, 0x9E))
        txt(s, contact, 3.2, 3.95, 6.9, 1.1,
            size=17, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    add_logo(s, logo_path, x=5.5, y=5.5, w=2.3, h=1.4)

    txt(s, 'SALES KIT  ·  POWERED BY AI', 0, SH - 0.5, SW, 0.45,
        size=10, color=RGBColor(0x4A, 0x7A, 0xD4), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# 投影片分派與主函式
# ══════════════════════════════════════════

def build_slide(prs, slide_data: dict, logo_path):
    """依據 type 欄位呼叫對應的投影片建立函式"""
    t = slide_data.get('type', 'content')
    if   t == 'cover':    slide_cover(prs, slide_data, logo_path)
    elif t == 'features': slide_features(prs, slide_data)
    elif t == 'stats':    slide_stats(prs, slide_data)
    elif t == 'steps':    slide_steps(prs, slide_data)
    elif t == 'case':     slide_case(prs, slide_data)       # BNI 成功案例
    elif t == 'cta':      slide_cta(prs, slide_data, logo_path)
    else:                 slide_content(prs, slide_data)  # intro / problem / solution


def generate_sales_kit(form_data: dict, output_path: str):
    """
    主進入點：
    1. 呼叫 Claude API 取得投影片 JSON
    2. 建立 Presentation 物件
    3. 逐頁生成投影片
    4. 儲存 .pptx 檔案
    """
    # Step 1: 取得 AI 生成的內容
    result = call_claude(form_data)
    slides = result.get('slides', [])
    if not slides:
        raise ValueError("AI 未回傳投影片資料，請再試一次")

    # Step 2: 初始化簡報（16:9 寬螢幕）
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    # Step 3: 逐頁建立投影片
    logo_path = form_data.get('logo_path')
    for slide_data in slides:
        build_slide(prs, slide_data, logo_path)

    # Step 4: 儲存
    prs.save(output_path)
